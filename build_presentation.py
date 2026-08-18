import pptx
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import shutil

def create_presentation():
    prs = pptx.Presentation('template_presentation.pptx')
    
    # Elegant Color Palette
    NAVY = RGBColor(15, 23, 42)        # #0F172A (Headings)
    SLATE = RGBColor(51, 65, 85)       # #334155 (Body text)
    DARK_BLUE = RGBColor(29, 78, 216)  # #1D4ED8 (Primary accent)
    TEAL = RGBColor(13, 148, 136)      # #0D9488 (Secondary accent)
    ORANGE = RGBColor(217, 119, 6)     # #D97706 (Bullet / Highlight)
    
    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.name == 'Google Shape;91;p1':
            shape.left = Emu(450000)
            shape.top = Emu(2150000)
            shape.width = Emu(11300000)
            shape.height = Emu(3150000)
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            fields = [
                ('Project Name (Registered on portal): ', 'Smart Nagpur'),
                ('Registration No (Received on Whatsapp): ', '[Your Registration No / Team ID]'),
                ('Theme (Existing): ', 'Smart Cities, Urban Governance & Clean Environment'),
                ('Problem Statement Title (Existing): ', 'Digital Civic Grievance Redressal, Street Vendor Permitting & Municipal Maintenance'),
                ('Expected Solution Title (Existing): ', 'Centralized Civic Service Delivery & Municipal Workflow Automation Platform'),
                ('Designed Solution Title (Team): ', 'Smart Nagpur: Multi-Role Civic Ecosystem (Citizen, Staff, Admin) with Real-Time Geotagged Tracking & Supabase RLS')
            ]
            
            for i, (label, val) in enumerate(fields):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(7)
                
                r1 = p.add_run()
                r1.text = label
                r1.font.name = 'Arial'
                r1.font.size = Pt(13)
                r1.font.bold = True
                r1.font.color.rgb = NAVY
                
                r2 = p.add_run()
                r2.text = val
                r2.font.name = 'Arial'
                r2.font.size = Pt(13)
                r2.font.bold = (i == 0 or i == 5)
                r2.font.color.rgb = DARK_BLUE if (i == 0 or i == 5) else SLATE

    # ----------------------------------------------------
    # SLIDE 2: Solution
    # ----------------------------------------------------
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if shape.name == 'Google Shape;104;p2': # PROJECT TITLE
            shape.left = Emu(2000000)
            shape.top = Emu(1400000)
            shape.width = Emu(8000000)
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = 'SMART NAGPUR: UNIFIED CIVIC ENGAGEMENT PLATFORM'
            r.font.name = 'Arial'
            r.font.size = Pt(17)
            r.font.bold = True
            r.font.color.rgb = DARK_BLUE
            
        elif shape.name == 'Google Shape;101;p2': # Solution content
            shape.left = Emu(500000)
            shape.top = Emu(2050000)
            shape.width = Emu(11200000)
            shape.height = Emu(4150000)
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            points = [
                ('1. 3-Tier Multi-Role Mobile Ecosystem: ', 'Built with Flutter for 3 dedicated user personas -- Citizen App (com.smartnagpur.citizen), Field Staff App (com.smartnagpur.staff), and Municipal Admin Portal (com.smartnagpur.admin).'),
                ('2. 10 Core Municipal Service Categories: ', 'Comprehensive grievance reporting covering Garbage/Sanitation, Water Supply, Roads & Potholes, Drainage, Streetlights, Stray Animals, Encroachments, Public Spaces, Vendor Permitting, and General Inquiries.'),
                ('3. 4-Step Grievance Reporting Wizard: ', 'Instant GPS auto-fetching, interactive map pin refinement, multi-image proof capture (up to 3 images), and atomic cloud submission via Supabase.'),
                ('4. Digital Street Vendor Permitting: ', 'End-to-end 4-step digital onboarding, document verification (Aadhaar/FSSAI), designated hawking zone explorer, and express permit renewals.'),
                ('5. Real-Time SLA Transparency & Milestones: ', 'Live milestone audit trail (Submitted -> Under Review -> Assigned -> In Progress -> Resolved) with official officer remarks and native Marathi/English localization.'),
                ('6. Innovation & Uniqueness: ', 'Single codebase supporting 3 build flavors, offline demo sandbox, transactional database RPCs, and zero-latency real-time field task dispatching.')
            ]
            
            for i, (bold_prefix, text_body) in enumerate(points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(6)
                
                r_bullet = p.add_run()
                r_bullet.text = '> '
                r_bullet.font.name = 'Arial'
                r_bullet.font.size = Pt(12)
                r_bullet.font.bold = True
                r_bullet.font.color.rgb = ORANGE
                
                r_head = p.add_run()
                r_head.text = bold_prefix
                r_head.font.name = 'Arial'
                r_head.font.size = Pt(12)
                r_head.font.bold = True
                r_head.font.color.rgb = NAVY
                
                r_body = p.add_run()
                r_body.text = text_body
                r_body.font.name = 'Arial'
                r_body.font.size = Pt(12)
                r_body.font.bold = False
                r_body.font.color.rgb = SLATE

    # ----------------------------------------------------
    # SLIDE 3: Technical Approach
    # ----------------------------------------------------
    s3 = prs.slides[2]
    for shape in s3.shapes:
        if shape.name == 'Google Shape;113;p3':
            shape.left = Emu(500000)
            shape.top = Emu(2100000)
            shape.width = Emu(11200000)
            shape.height = Emu(4150000)
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            sections = [
                ('A. TECHNOLOGY STACK & ARCHITECTURE', [
                    ('Frontend / Mobile: ', 'Flutter 3.x (Dart), Riverpod for reactive state management, Material 3 Design System, Marathi & English localization.'),
                    ('Backend & Database: ', 'Supabase (PostgreSQL 15), Row-Level Security (RLS) policies, Custom transactional RPCs (submit_complaint_atomic, assign_staff_task).'),
                    ('Auth & Storage: ', 'Supabase Auth (JWT/RBAC), Encrypted Storage Buckets (complaint-photos, vendor-docs) with private access control.'),
                    ('Geolocation & Maps: ', 'geolocator GPS accuracy engine with OpenStreetMap / Mapbox pin rendering.')
                ]),
                ('B. IMPLEMENTATION WORKFLOW & PROCESS', [
                    ('Citizen Submission: ', 'Grievance captured with GPS coordinates & photo evidence -> atomic database insert.'),
                    ('Admin Triage & RBAC: ', 'Municipal officers review queue, verify documents, set SLA priority, and dispatch to on-duty field staff.'),
                    ('Field Staff Execution: ', 'Real-time task notification -> GPS route navigation -> on-ground resolution -> Before/After photo proof upload -> instant ticket closure.')
                ])
            ]
            
            p_first = True
            for sec_title, items in sections:
                p_sec = tf.paragraphs[0] if p_first else tf.add_paragraph()
                p_first = False
                p_sec.space_before = Pt(3)
                p_sec.space_after = Pt(2)
                
                r_sec = p_sec.add_run()
                r_sec.text = sec_title
                r_sec.font.name = 'Arial'
                r_sec.font.size = Pt(12)
                r_sec.font.bold = True
                r_sec.font.color.rgb = DARK_BLUE
                
                for prefix, body in items:
                    p = tf.add_paragraph()
                    p.space_after = Pt(4)
                    
                    r_dot = p.add_run()
                    r_dot.text = '  > '
                    r_dot.font.name = 'Arial'
                    r_dot.font.size = Pt(11)
                    r_dot.font.bold = True
                    r_dot.font.color.rgb = ORANGE
                    
                    r_pre = p.add_run()
                    r_pre.text = prefix
                    r_pre.font.name = 'Arial'
                    r_pre.font.size = Pt(11)
                    r_pre.font.bold = True
                    r_pre.font.color.rgb = NAVY
                    
                    r_b = p.add_run()
                    r_b.text = body
                    r_b.font.name = 'Arial'
                    r_b.font.size = Pt(11)
                    r_b.font.bold = False
                    r_b.font.color.rgb = SLATE

    # ----------------------------------------------------
    # SLIDE 4: Practical Implementation
    # ----------------------------------------------------
    s4 = prs.slides[3]
    for shape in s4.shapes:
        if shape.name == 'Google Shape;123;p4':
            shape.left = Emu(500000)
            shape.top = Emu(2100000)
            shape.width = Emu(11200000)
            shape.height = Emu(4150000)
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            sections = [
                ('1. PRACTICAL ANALYSIS & IMPLEMENTATION', [
                    ('Production Codebase: ', 'Unified Flutter project configured with 3 build flavors (Citizen, Staff, Admin) sharing domain models and repositories.'),
                    ('Enterprise Security: ', 'PostgreSQL Row-Level Security (RLS) protects citizen privacy, isolates vendor documentation, and restricts admin operations.')
                ]),
                ('2. FORESEEN CHALLENGES', [
                    ('Digital Literacy: ', 'Street vendors and elderly citizens may find complex digital interfaces challenging.'),
                    ('Spam / Duplicate Reports: ', 'Multiple citizens reporting the same civic issue (e.g. single pothole) causing clutter.'),
                    ('Field Connectivity: ', 'Intermittent network connectivity in remote wards and basements.')
                ]),
                ('3. STRATEGIES FOR OVERCOMING CHALLENGES', [
                    ('Accessible UI & Marathi Support: ', 'Clean, icon-driven interface with native Marathi localization and voice-assisted form inputs.'),
                    ('Geospatial De-duplication: ', 'Automatic grouping of nearby complaints within 50m radius into a single municipal ticket.'),
                    ('Offline-First Architecture: ', 'Local database caching with automatic background synchronization upon reconnection.')
                ])
            ]
            
            p_first = True
            for sec_title, items in sections:
                p_sec = tf.paragraphs[0] if p_first else tf.add_paragraph()
                p_first = False
                p_sec.space_before = Pt(3)
                p_sec.space_after = Pt(2)
                
                r_sec = p_sec.add_run()
                r_sec.text = sec_title
                r_sec.font.name = 'Arial'
                r_sec.font.size = Pt(12)
                r_sec.font.bold = True
                r_sec.font.color.rgb = DARK_BLUE
                
                for prefix, body in items:
                    p = tf.add_paragraph()
                    p.space_after = Pt(3)
                    
                    r_dot = p.add_run()
                    r_dot.text = '  > '
                    r_dot.font.name = 'Arial'
                    r_dot.font.size = Pt(11)
                    r_dot.font.bold = True
                    r_dot.font.color.rgb = ORANGE
                    
                    r_pre = p.add_run()
                    r_pre.text = prefix
                    r_pre.font.name = 'Arial'
                    r_pre.font.size = Pt(11)
                    r_pre.font.bold = True
                    r_pre.font.color.rgb = NAVY
                    
                    r_b = p.add_run()
                    r_b.text = body
                    r_b.font.name = 'Arial'
                    r_b.font.size = Pt(11)
                    r_b.font.bold = False
                    r_b.font.color.rgb = SLATE

    # ----------------------------------------------------
    # SLIDE 5: Impact and Benefits
    # ----------------------------------------------------
    s5 = prs.slides[4]
    for shape in s5.shapes:
        if shape.name == 'Google Shape;134;p5':
            shape.left = Emu(500000)
            shape.top = Emu(2100000)
            shape.width = Emu(11200000)
            shape.height = Emu(4150000)
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            sections = [
                ('1. POTENTIAL IMPACT & STAKEHOLDER BENEFITS', [
                    ('Citizens of Nagpur: ', 'Grievance resolution turnaround cut from weeks to 24-48 hours; complete transparency with before/after photos and audit logs.'),
                    ('Nagpur Municipal Corporation (NMC): ', '70% reduction in paper processing; centralized live operations dashboard with ward-wise analytics and SLA heatmaps.'),
                    ('Street Vendors & Hawkers: ', 'Transparent digital permits, elimination of unofficial middlemen, and legal protection in designated hawking zones.'),
                    ('Field Maintenance Staff: ', 'Optimized task routing, clear daily workload visibility, and verifiable digital proof of work completion.')
                ]),
                ('2. ALIGNMENT WITH SUSTAINABLE DEVELOPMENT GOALS (SDGs)', [
                    ('SDG 11 (Sustainable Cities & Communities): ', 'Targeted improvement in urban waste management, pothole-free roads, and public safety infrastructure.'),
                    ('SDG 9 (Industry, Innovation & Infrastructure): ', 'Modern digital infrastructure elevating public municipal utility services.'),
                    ('SDG 16 (Peace, Justice & Strong Institutions): ', 'Transparent, accountable, corruption-free public grievance redressal with verified audit trails.')
                ])
            ]
            
            p_first = True
            for sec_title, items in sections:
                p_sec = tf.paragraphs[0] if p_first else tf.add_paragraph()
                p_first = False
                p_sec.space_before = Pt(3)
                p_sec.space_after = Pt(2)
                
                r_sec = p_sec.add_run()
                r_sec.text = sec_title
                r_sec.font.name = 'Arial'
                r_sec.font.size = Pt(12)
                r_sec.font.bold = True
                r_sec.font.color.rgb = DARK_BLUE
                
                for prefix, body in items:
                    p = tf.add_paragraph()
                    p.space_after = Pt(3)
                    
                    r_dot = p.add_run()
                    r_dot.text = '  > '
                    r_dot.font.name = 'Arial'
                    r_dot.font.size = Pt(11)
                    r_dot.font.bold = True
                    r_dot.font.color.rgb = ORANGE
                    
                    r_pre = p.add_run()
                    r_pre.text = prefix
                    r_pre.font.name = 'Arial'
                    r_pre.font.size = Pt(11)
                    r_pre.font.bold = True
                    r_pre.font.color.rgb = NAVY
                    
                    r_b = p.add_run()
                    r_b.text = body
                    r_b.font.name = 'Arial'
                    r_b.font.size = Pt(11)
                    r_b.font.bold = False
                    r_b.font.color.rgb = SLATE

    # ----------------------------------------------------
    # SLIDE 6: Research and References
    # ----------------------------------------------------
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if shape.name == 'Google Shape;145;p6':
            shape.left = Emu(500000)
            shape.top = Emu(2100000)
            shape.width = Emu(11200000)
            shape.height = Emu(4150000)
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            sections = [
                ('1. MUNICIPAL & GOVERNANCE FRAMEWORKS', [
                    ('Nagpur Municipal Corporation (NMC) Citizen Charter: ', 'Standard Operating Procedures (SOPs) for civic grievance classification, escalation levels, and ward boundary definitions.'),
                    ('Smart Cities Mission Guidelines (MoHUA, Govt. of India): ', 'Standards for citizen-centric mobile governance and Integrated Command & Control Centre (ICCC) data models.'),
                    ('Street Vendors (Protection of Livelihood & Regulation of Street Vending) Act 2014: ', 'Legal guidelines for designated vending zones, survey mechanisms, and digital registration.')
                ]),
                ('2. TECHNICAL & SYSTEM ARCHITECTURE REFERENCES', [
                    ('Supabase Architecture & PostgreSQL Best Practices: ', 'Implementation of Row-Level Security (RLS), ACID-compliant transactional RPCs, and Realtime WebSocket replication.'),
                    ('Flutter Multi-Flavor Architecture: ', 'Industry-standard multi-flavor build separation for multi-role mobile deployments from a single maintainable codebase.')
                ])
            ]
            
            p_first = True
            for sec_title, items in sections:
                p_sec = tf.paragraphs[0] if p_first else tf.add_paragraph()
                p_first = False
                p_sec.space_before = Pt(4)
                p_sec.space_after = Pt(3)
                
                r_sec = p_sec.add_run()
                r_sec.text = sec_title
                r_sec.font.name = 'Arial'
                r_sec.font.size = Pt(12)
                r_sec.font.bold = True
                r_sec.font.color.rgb = DARK_BLUE
                
                for prefix, body in items:
                    p = tf.add_paragraph()
                    p.space_after = Pt(4)
                    
                    r_dot = p.add_run()
                    r_dot.text = '  > '
                    r_dot.font.name = 'Arial'
                    r_dot.font.size = Pt(11)
                    r_dot.font.bold = True
                    r_dot.font.color.rgb = ORANGE
                    
                    r_pre = p.add_run()
                    r_pre.text = prefix
                    r_pre.font.name = 'Arial'
                    r_pre.font.size = Pt(11)
                    r_pre.font.bold = True
                    r_pre.font.color.rgb = NAVY
                    
                    r_b = p.add_run()
                    r_b.text = body
                    r_b.font.name = 'Arial'
                    r_b.font.size = Pt(11)
                    r_b.font.bold = False
                    r_b.font.color.rgb = SLATE

    output_filename = 'SmartNagpur_Vikasit_Nagpur_Hackathon_2026.pptx'
    prs.save(output_filename)
    prs.save('template_presentation.pptx')
    
    # Also copy to Downloads folder so the user has it ready
    try:
        shutil.copyfile(output_filename, r'C:\Users\Shivam\Downloads\SmartNagpur_Vikasit_Nagpur_Hackathon_2026.pptx')
        print(f'Successfully copied to Downloads folder as well!')
    except Exception as e:
        print(f'Note: Could not copy to Downloads directly: {e}')
        
    print(f'Successfully generated {output_filename}!')

if __name__ == '__main__':
    create_presentation()
